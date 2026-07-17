#!/usr/bin/env python3
import re
import sys
from pathlib import Path

from PIL import Image as PILImage


def natural_key(path):
    return [int(part) if part.isdigit() else part for part in re.split(r"(\d+)", path.name)]


def page_images(directory):
    images = sorted(Path(directory).glob("page-*.png"), key=natural_key)
    if not images:
        raise RuntimeError("The PDF renderer did not create any page images.")
    return images


def build_docx(images, output):
    from docx import Document
    from docx.shared import Inches

    document = Document()
    section = document.sections[0]
    section.page_width = Inches(8.27)
    section.page_height = Inches(11.69)
    section.top_margin = Inches(0.25)
    section.bottom_margin = Inches(0.25)
    section.left_margin = Inches(0.25)
    section.right_margin = Inches(0.25)
    usable_width = section.page_width - section.left_margin - section.right_margin
    usable_height = section.page_height - section.top_margin - section.bottom_margin
    for index, image_path in enumerate(images):
        with PILImage.open(image_path) as image:
            ratio = image.width / max(1, image.height)
        width = usable_width
        height = int(width / ratio)
        if height > usable_height:
            height = usable_height
            width = int(height * ratio)
        paragraph = document.add_paragraph()
        paragraph.paragraph_format.space_before = 0
        paragraph.paragraph_format.space_after = 0
        paragraph.alignment = 1
        paragraph.add_run().add_picture(str(image_path), width=width, height=height)
        if index < len(images) - 1:
            document.add_page_break()
    document.save(output)


def build_pptx(images, output):
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]
    for image_path in images:
        slide = presentation.slides.add_slide(blank)
        with PILImage.open(image_path) as image:
            ratio = image.width / max(1, image.height)
        slide_ratio = presentation.slide_width / presentation.slide_height
        if ratio >= slide_ratio:
            width = presentation.slide_width
            height = int(width / ratio)
            left = 0
            top = int((presentation.slide_height - height) / 2)
        else:
            height = presentation.slide_height
            width = int(height * ratio)
            top = 0
            left = int((presentation.slide_width - width) / 2)
        slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
    presentation.save(output)


def build_xlsx(images, output):
    from openpyxl import Workbook
    from openpyxl.drawing.image import Image as WorksheetImage

    workbook = Workbook()
    for index, image_path in enumerate(images):
        worksheet = workbook.active if index == 0 else workbook.create_sheet()
        worksheet.title = f"Page {index + 1}"
        worksheet.sheet_view.showGridLines = False
        picture = WorksheetImage(str(image_path))
        if picture.width > 1000:
            scale = 1000 / picture.width
            picture.width = int(picture.width * scale)
            picture.height = int(picture.height * scale)
        worksheet.add_image(picture, "A1")
    workbook.save(output)


def main():
    if len(sys.argv) != 4:
        raise RuntimeError("Expected output format, output path, and page-image directory.")
    output_format, output_path, image_directory = sys.argv[1:]
    images = page_images(image_directory)
    if output_format == "docx":
        build_docx(images, output_path)
    elif output_format == "pptx":
        build_pptx(images, output_path)
    elif output_format == "xlsx":
        build_xlsx(images, output_path)
    else:
        raise RuntimeError(f"Unsupported generated Office format: {output_format}")


if __name__ == "__main__":
    main()
